import tempfile
import subprocess
import uuid
import os
from pathlib import Path
import datetime
import shutil
import traceback
from fastapi import HTTPException
from fastapi.responses import JSONResponse, FileResponse
from pathlib import Path
import base64
from io import BytesIO
import fitz
from PIL import Image
import json
import hashlib
from functools import lru_cache
import docx
import openpyxl
import pptx
from fastapi.responses import HTMLResponse

from routers.utils.misc_files_utils import *
from routers.utils.misc_keycloak_utils import *


# Cache for PDF documents to avoid reopening frequently
pdf_cache = {}
CACHE_SIZE_LIMIT = 10  # Maximum number of PDFs to keep in memory

def get_pdf_cache_key(abs_path, mtime):
    """Generate a cache key for PDF based on path and modification time"""
    return f"{abs_path}_{mtime}"

def get_cached_pdf(abs_path):
    """Get PDF from cache or load it"""
    try:
        mtime = os.path.getmtime(abs_path)
        cache_key = get_pdf_cache_key(abs_path, mtime)
        
        # Clean old cache entries if cache is too large
        if len(pdf_cache) > CACHE_SIZE_LIMIT:
            # Remove oldest entries
            oldest_keys = list(pdf_cache.keys())[:-CACHE_SIZE_LIMIT//2]
            for old_key in oldest_keys:
                if old_key in pdf_cache:
                    pdf_cache[old_key].close()
                    del pdf_cache[old_key]
        
        # Check if PDF is already cached
        if cache_key in pdf_cache:
            return pdf_cache[cache_key]
        
        # Load new PDF
        doc = fitz.open(abs_path)
        pdf_cache[cache_key] = doc
        return doc
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error loading PDF: {str(e)}")


async def file_preview(path):
    base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
    relative_path = path.lstrip("/\\")
    abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
    # print(abs_path)
    if not os.path.isfile(abs_path):
        raise HTTPException(status_code=400, detail="Invalid file path")
    ext = os.path.splitext(abs_path)[1].lower()
    try:
        if ext == ".pdf":
            doc = fitz.open(abs_path)
            if doc.page_count < 1:
                raise HTTPException(status_code=500, detail="Could not render PDF")
            page = doc.load_page(0)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img_data = pix.tobytes("png")
            img = Image.open(BytesIO(img_data))
        elif ext in [".png", ".jpg", ".jpeg"]:
            img = Image.open(abs_path)
        else:
            raise HTTPException(status_code=415, detail="Preview not supported for this file type")

        img.thumbnail((100, 100))
        buffered = BytesIO()
        img.save(buffered, format="PNG")
        preview_path = os.path.join(os.getcwd(), "preview/preview_output.png")
        img.save(preview_path, format="PNG") # save preview
        img_b64 = base64.b64encode(buffered.getvalue()).decode("utf-8")
        return img_b64
        
    except Exception as e:
        raise e from e
    

async def download_file(path: str, user_id: str = None, username: str = None):
    try:
        base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
        relative_path = path.lstrip("/\\")
        abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
        download_file_path = Path(abs_path).resolve()
        if not download_file_path.is_file():
            raise HTTPException(status_code=404, detail=f"File not found: {download_file_path}")
        
        # Update user's recent_files attribute in Keycloak if user info is provided
        if user_id and username:
            await update_user_recent_file_attribute(user_id, username, path)
        
        return FileResponse(
        path=download_file_path,
        filename=download_file_path.name
        )
    except Exception as e:
        raise e from e
    

async def search_files(search_str: str):
    try:
        base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
        results = search_files_and_folders(base_dir, search_str)
        return results
    except Exception as e:
        raise e


async def dir_contents(path: str, permissions: list, roles: list):
    try:
        base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
        relative_path = path.lstrip("/\\")
        abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
        p_abs_path = Path(abs_path)
        if not p_abs_path.is_dir(): raise HTTPException(status_code=404, detail="path does not exist")
        print("permissions:", permissions)
        results = await dir_contents_details(abs_path, permissions, roles)
        return results

    except Exception as e:
        raise e from e


async def create_dir(path: str):
    try:
        base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
        relative_path = str(Path(path.lstrip("/\\")).as_posix())
        abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
        os.makedirs(abs_path, exist_ok=False)
        
        # resource_payload = (                            
        #     {
        #         "name":relative_path,
        #         "displayName":relative_path,
        #         "type":"dir",
        #         "icon_uri":"",
        #         "ownerManagedAccess":False,
        #         "attributes":{},
        #         "scopes":[]
        #     }
        # )
        # await create_resource(resource_payload)
    
        return relative_path
    except Exception as e:
        raise e


async def delete_file_and_dir(path: str):
    try:
        base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
        relative_path = str(Path(path.lstrip("/\\")).as_posix())
        # print("relative_path:", relative_path)
        abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
        all_resources = await get_all_resources()
        # print("all_resources:", all_resources)
        if os.path.isfile(abs_path):
            os.remove(abs_path)
        elif os.path.isdir(abs_path):
            shutil.rmtree(abs_path)
        else:
            raise HTTPException(status_code=400, detail="invalid path")
        for resource in all_resources:
                if relative_path in resource: await delete_resource(all_resources[resource])
        return f"deleted: {relative_path}"
    except Exception as e:
        raise e from e
    

async def upload_files(folder: str, files, path: str):
    try:
        base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
        relative_path = path.lstrip("/\\") if path else ""
        if folder:
            relative_path = os.path.join(relative_path, folder)
            # print("relative_path:", relative_path)
            await create_dir(relative_path)
        abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
        uploaded_files = []
        for file in files:
            file_path = os.path.join(abs_path, file.filename)
            with open(file_path, "wb") as buffer:
                shutil.copyfileobj(file.file, buffer)
            relative_file_location = str(Path(os.path.join(relative_path, file.filename)).as_posix()) # e.g. docs/test_file.jpg
            # resource_payload = (
            #     {
            #         "name":relative_file_location,
            #         "displayName":relative_file_location,
            #         "type":"file",
            #         "icon_uri":"",
            #         "ownerManagedAccess":False,
            #         "attributes":{},
            #         "scopes":[]
            #     }
            # )
            # await create_resource(resource_payload)

            uploaded_files.append(file.filename)

        return uploaded_files
    
    except Exception as e:
        tb_str = traceback.format_exc()
        print(f"Error in upload_files: {tb_str}")
        raise e from e


async def upload_multiple_folders(files, directory_structure_json: str):
    """
    Upload multiple folders with complex directory structures.
    
    Args:
        files: List of uploaded files
        directory_structure_json: Serialized JSON containing the directory structure
        
    Returns:
        Dictionary with upload results
    """
    try:
        base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
        
        # Parse the directory structure
        directory_structure = json.loads(directory_structure_json)
        print("directory_structure:", directory_structure)
          # Create a mapping of file names to file objects for quick lookup
        file_map = {file.filename: file for file in files}
        
        uploaded_files = []
        created_dirs = []
        
        # Recursively process the directory structure
        await process_directory_structure(
            directory_structure, 
            base_dir, 
            "", 
            file_map, 
            uploaded_files, 
            created_dirs
        )
        
        return {
            "uploaded_files": uploaded_files,
            "created_directories": created_dirs,
            "total_files": len(uploaded_files),
            "total_directories": len(created_dirs)
        }
        
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=400, detail=f"Invalid JSON in directory_structure: {str(e)}")
    except Exception as e:
        tb_str = traceback.format_exc()
        print(f"Error in upload_multiple_folders: {tb_str}")
        raise HTTPException(status_code=500, detail=f"Upload failed: {str(e)}")


async def get_pdf_info(path):
    """Get PDF metadata including page count, dimensions, etc."""
    base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
    relative_path = path.lstrip("/\\")
    abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
    
    if not os.path.isfile(abs_path):
        raise HTTPException(status_code=400, detail="Invalid file path")
    
    ext = os.path.splitext(abs_path)[1].lower()
    if ext != ".pdf":
        raise HTTPException(status_code=415, detail="File is not a PDF")
    
    try:
        doc = get_cached_pdf(abs_path)
        
        # Get first page to determine dimensions
        first_page = doc.load_page(0)
        rect = first_page.rect
        
        # Extract basic metadata
        metadata = doc.metadata
        
        pdf_info = {
            "page_count": doc.page_count,
            "width": rect.width,
            "height": rect.height,
            "title": metadata.get("title", ""),
            "author": metadata.get("author", ""),
            "subject": metadata.get("subject", ""),
            "creator": metadata.get("creator", ""),
            "producer": metadata.get("producer", ""),
            "creation_date": metadata.get("creationDate", ""),
            "modification_date": metadata.get("modDate", ""),
            "file_size": os.path.getsize(abs_path)
        }
        
        return pdf_info
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading PDF info: {str(e)}")


async def get_pdf_page(path, page_num, quality="medium", scale=1.0):
    """Get a specific page from PDF as base64 image"""
    base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
    relative_path = path.lstrip("/\\")
    abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
    
    if not os.path.isfile(abs_path):
        raise HTTPException(status_code=400, detail="Invalid file path")
    
    try:
        doc = get_cached_pdf(abs_path)
        
        if page_num < 1 or page_num > doc.page_count:
            raise HTTPException(status_code=400, detail=f"Page {page_num} not found. PDF has {doc.page_count} pages")
        
        # Load the specific page (0-indexed)
        page = doc.load_page(page_num - 1)
        
        # Set quality based on parameter
        quality_settings = {
            "low": 1.0,
            "medium": 1.5,
            "high": 2.0
        }
        base_scale = quality_settings.get(quality, 1.5)
        final_scale = base_scale * scale
        
        # Render page to image
        matrix = fitz.Matrix(final_scale, final_scale)
        pix = page.get_pixmap(matrix=matrix)
        img_data = pix.tobytes("png")
        
        # Convert to PIL Image for additional processing if needed
        img = Image.open(BytesIO(img_data))
        
        # Convert to base64
        buffered = BytesIO()
        img.save(buffered, format="PNG", optimize=True)
        img_b64 = base64.b64encode(buffered.getvalue()).decode("utf-8")
        
        return {
            "page_number": page_num,
            "image_data": img_b64,
            "width": img.width,
            "height": img.height,
            "scale": final_scale
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error rendering PDF page: {str(e)}")


async def get_pdf_pages_range(path, start_page, end_page, quality="medium", scale=1.0):
    """Get multiple PDF pages in a range"""
    base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
    relative_path = path.lstrip("/\\")
    abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
    
    if not os.path.isfile(abs_path):
        raise HTTPException(status_code=400, detail="Invalid file path")
    
    try:
        doc = get_cached_pdf(abs_path)
        
        # Validate page range
        if start_page < 1 or end_page > doc.page_count or start_page > end_page:
            raise HTTPException(status_code=400, detail=f"Invalid page range. PDF has {doc.page_count} pages")
        
        # Limit range to prevent memory issues
        max_pages_per_request = 5
        if end_page - start_page + 1 > max_pages_per_request:
            raise HTTPException(status_code=400, detail=f"Maximum {max_pages_per_request} pages per request")
        
        pages_data = []
        
        for page_num in range(start_page, end_page + 1):
            page_data = await get_pdf_page(path, page_num, quality, scale)
            pages_data.append(page_data)
        
        return {
            "pages": pages_data,
            "start_page": start_page,
            "end_page": end_page,
            "total_pages": len(pages_data)
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error rendering PDF pages: {str(e)}")


async def search_pdf_text(path, search_text):
    """Search for text within PDF and return page numbers and positions"""
    base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
    relative_path = path.lstrip("/\\")
    abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
    
    if not os.path.isfile(abs_path):
        raise HTTPException(status_code=400, detail="Invalid file path")
    
    if not search_text or len(search_text.strip()) < 2:
        raise HTTPException(status_code=400, detail="Search text must be at least 2 characters")
    
    try:
        doc = get_cached_pdf(abs_path)
        search_results = []
        
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            
            # Search for text on the page
            text_instances = page.search_for(search_text)
            
            if text_instances:
                page_results = {
                    "page_number": page_num + 1,
                    "matches": []
                }
                
                for rect in text_instances:
                    # Get surrounding context
                    words = page.get_text("words")
                    context = ""
                    for word in words:
                        word_rect = fitz.Rect(word[:4])
                        if word_rect.intersects(rect):
                            # Get some words before and after for context
                            word_index = words.index(word)
                            start_idx = max(0, word_index - 5)
                            end_idx = min(len(words), word_index + 6)
                            context_words = [w[4] for w in words[start_idx:end_idx]]
                            context = " ".join(context_words)
                            break
                    
                    page_results["matches"].append({
                        "position": {
                            "x": rect.x0,
                            "y": rect.y0,
                            "width": rect.width,
                            "height": rect.height
                        },
                        "context": context[:200]  # Limit context length
                    })
                
                search_results.append(page_results)
        
        return {
            "search_text": search_text,
            "total_matches": sum(len(page["matches"]) for page in search_results),
            "pages_with_matches": len(search_results),
            "results": search_results
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error searching PDF: {str(e)}")


async def get_pdf_page_with_text(path, page_num, quality="medium", scale=1.0):
    """Get a PDF page with both image and text layer data"""
    base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
    relative_path = path.lstrip("/\\")
    abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
    
    if not os.path.isfile(abs_path):
        raise HTTPException(status_code=400, detail="Invalid file path")
    
    try:
        doc = get_cached_pdf(abs_path)
        
        if page_num < 1 or page_num > doc.page_count:
            raise HTTPException(status_code=400, detail=f"Page {page_num} not found")
        
        page = doc.load_page(page_num - 1)
        
        # Get the image data (same as before)
        quality_settings = {"low": 1.0, "medium": 1.5, "high": 2.0}
        base_scale = quality_settings.get(quality, 1.5)
        final_scale = base_scale * scale
        
        matrix = fitz.Matrix(final_scale, final_scale)
        pix = page.get_pixmap(matrix=matrix)
        img_data = pix.tobytes("png")
        
        img = Image.open(BytesIO(img_data))
        buffered = BytesIO()
        img.save(buffered, format="PNG", optimize=True)
        img_b64 = base64.b64encode(buffered.getvalue()).decode("utf-8")
        
        # Get text layer data
        text_data = await get_pdf_text_layer(path, page_num, final_scale)
        
        return {
            "page_number": page_num,
            "image_data": img_b64,
            "width": img.width,
            "height": img.height,
            "scale": final_scale,
            "text_layer": text_data["text_blocks"]
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error rendering PDF page with text: {str(e)}")


async def get_pdf_text_layer(path, page_num, scale=1.0):
    """Get text layer data for a PDF page with positioning"""
    base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
    relative_path = path.lstrip("/\\")
    abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
    
    if not os.path.isfile(abs_path):
        raise HTTPException(status_code=400, detail="Invalid file path")
    
    try:
        doc = get_cached_pdf(abs_path)
        
        if page_num < 1 or page_num > doc.page_count:
            raise HTTPException(status_code=400, detail=f"Page {page_num} not found")
        
        page = doc.load_page(page_num - 1)
        
        # Get text blocks with positioning
        text_blocks = []
        
        # Method 1: Get text with detailed positioning (words)
        words = page.get_text("words")
        for word in words:
            x0, y0, x1, y1, text, block_no, line_no, word_no = word
            
            # Scale coordinates to match the rendered image
            text_blocks.append({
                "text": text,
                "bbox": {
                    "x": x0 * scale,
                    "y": y0 * scale,
                    "width": (x1 - x0) * scale,
                    "height": (y1 - y0) * scale
                },
                "block_no": block_no,
                "line_no": line_no,
                "word_no": word_no,
                "type": "word"
            })
        
        # Method 2: Get text blocks (paragraphs) for better structure
        blocks = page.get_text("dict")["blocks"]
        text_paragraphs = []
        
        for block in blocks:
            if "lines" in block:  # Text block
                for line in block["lines"]:
                    line_text = ""
                    line_bbox = None
                    
                    for span in line["spans"]:
                        line_text += span["text"]
                        span_bbox = span["bbox"]
                        
                        if line_bbox is None:
                            line_bbox = list(span_bbox)
                        else:
                            # Extend bounding box
                            line_bbox[0] = min(line_bbox[0], span_bbox[0])
                            line_bbox[1] = min(line_bbox[1], span_bbox[1])
                            line_bbox[2] = max(line_bbox[2], span_bbox[2])
                            line_bbox[3] = max(line_bbox[3], span_bbox[3])
                    
                    if line_text.strip():
                        text_paragraphs.append({
                            "text": line_text,
                            "bbox": {
                                "x": line_bbox[0] * scale,
                                "y": line_bbox[1] * scale,
                                "width": (line_bbox[2] - line_bbox[0]) * scale,
                                "height": (line_bbox[3] - line_bbox[1]) * scale
                            },
                            "type": "line",
                            "font_info": {
                                "size": line["spans"][0].get("size", 12) * scale if line["spans"] else 12,
                                "font": line["spans"][0].get("font", "unknown") if line["spans"] else "unknown"
                            }
                        })
        
        return {
            "page_number": page_num,
            "scale": scale,
            "text_blocks": text_blocks,
            "text_paragraphs": text_paragraphs,
            "page_width": page.rect.width * scale,
            "page_height": page.rect.height * scale
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error extracting text layer: {str(e)}")


async def get_raw_pdf(path):
    """Serve raw PDF file for PDF.js viewer"""
    base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
    relative_path = path.lstrip("/\\")
    abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
    
    if not os.path.isfile(abs_path):
        raise HTTPException(status_code=404, detail="File not found")
    
    ext = os.path.splitext(abs_path)[1].lower()
    if ext != ".pdf":
        raise HTTPException(status_code=415, detail="File is not a PDF")
    
    try:
        return FileResponse(
            path=abs_path,
            media_type="application/pdf",
            filename=os.path.basename(abs_path),
            headers={
                "Cache-Control": "public, max-age=3600",
                "Content-Disposition": "inline"  # Display in browser instead of download
            }
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error serving PDF: {str(e)}")


async def get_docx_info(path):
    base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
    relative_path = path.lstrip("/\\")
    abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
    if not os.path.isfile(abs_path):
        raise HTTPException(status_code=400, detail="Invalid file path")
    try:
        doc = docx.Document(abs_path)
        num_paragraphs = len(doc.paragraphs)
        num_tables = len(doc.tables)
        # Approximate 'pages' by splitting every 30 paragraphs
        page_count = max(1, (num_paragraphs + 29) // 30)
        return {
            "paragraphs": num_paragraphs,
            "tables": num_tables,
            "page_count": page_count,
            "file_size": os.path.getsize(abs_path)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading DOCX: {str(e)}")

async def get_docx_page(path, page_num):
    base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
    relative_path = path.lstrip("/\\")
    abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
    if not os.path.isfile(abs_path):
        raise HTTPException(status_code=400, detail="Invalid file path")
    temp_pdf = None
    try:
        # Find unoconv.exe in the current Python environment
        import sys
        venv_scripts = os.path.join(sys.prefix, 'Scripts')
        unoconv_path = os.path.join(venv_scripts, 'unoconv.exe')
        if not os.path.isfile(unoconv_path):
            # Try just 'unoconv' (for global install or Linux)
            unoconv_path = 'unoconv'

        # Place temp PDF in the same directory as the DOCX (base_dir)
        temp_pdf_name = f"docx2pdf_{uuid.uuid4().hex}.pdf"
        temp_pdf = os.path.normpath(os.path.join(base_dir, temp_pdf_name))
        # Use unoconv to convert DOCX to PDF (all absolute paths, same dir)
        result = subprocess.run([
            unoconv_path, "-f", "pdf", "-o", temp_pdf, abs_path
        ], capture_output=True)
        if result.returncode != 0 or not os.path.isfile(temp_pdf):
            raise HTTPException(status_code=500, detail=f"DOCX to PDF conversion failed: {result.stderr.decode('utf-8')}")
        # Use get_pdf_page to render the page
        from routers.utils.api_files_utils import get_pdf_page
        page_data = await get_pdf_page(temp_pdf, page_num)
        # Add a hint that this is a DOCX preview via PDF
        page_data["source"] = "docx->pdf"
        return page_data
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error rendering DOCX preview: {str(e)}")
    finally:
        # Clean up temp PDF
        if temp_pdf and os.path.isfile(temp_pdf):
            try:
                os.remove(temp_pdf)
            except Exception:
                pass

async def get_xlsx_info(path):
    base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
    relative_path = path.lstrip("/\\")
    abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
    if not os.path.isfile(abs_path):
        raise HTTPException(status_code=400, detail="Invalid file path")
    try:
        wb = openpyxl.load_workbook(abs_path, read_only=True)
        sheet_names = wb.sheetnames
        return {
            "sheet_names": sheet_names,
            "sheet_count": len(sheet_names),
            "file_size": os.path.getsize(abs_path)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading XLSX: {str(e)}")

async def get_xlsx_sheet(path, sheet_name):
    base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
    relative_path = path.lstrip("/\\")
    abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
    if not os.path.isfile(abs_path):
        raise HTTPException(status_code=400, detail="Invalid file path")
    try:
        wb = openpyxl.load_workbook(abs_path, read_only=True)
        if sheet_name not in wb.sheetnames:
            raise HTTPException(status_code=404, detail="Sheet not found")
        ws = wb[sheet_name]
        html = "<table border='1'>"
        for row in ws.iter_rows(values_only=True):
            html += "<tr>" + "".join(f"<td>{cell if cell is not None else ''}</td>" for cell in row) + "</tr>"
        html += "</table>"
        return {"sheet": sheet_name, "html": html}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading XLSX sheet: {str(e)}")

async def get_pptx_info(path):
    base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
    relative_path = path.lstrip("/\\")
    abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
    if not os.path.isfile(abs_path):
        raise HTTPException(status_code=400, detail="Invalid file path")
    try:
        from pptx import Presentation
        prs = Presentation(abs_path)
        slide_count = len(prs.slides)
        return {
            "slide_count": slide_count,
            "file_size": os.path.getsize(abs_path)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading PPTX: {str(e)}")

async def get_pptx_slide(path, slide_num):
    base_dir = os.path.normpath(os.path.join(os.getcwd(), "remote"))
    relative_path = path.lstrip("/\\")
    abs_path = os.path.normpath(os.path.join(base_dir, relative_path))
    if not os.path.isfile(abs_path):
        raise HTTPException(status_code=400, detail="Invalid file path")
    try:
        from pptx import Presentation
        prs = Presentation(abs_path)
        if slide_num < 1 or slide_num > len(prs.slides):
            raise HTTPException(status_code=404, detail="Slide not found")
        slide = prs.slides[slide_num - 1]
        # Render slide text as HTML (images not included)
        html = "<div>"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                html += f"<p>{shape.text}</p>"
        html += "</div>"
        return {"slide": slide_num, "html": html}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading PPTX slide: {str(e)}")